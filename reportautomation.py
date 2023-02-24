from configparser import Interpolation
from ctypes import util
from pickle import TRUE
from re import X
#from turtle import color
#from h11 import Data
from matplotlib import image
from pyparsing import col
import streamlit as st
import numpy as np
import pandas as pd
from wordcloud import STOPWORDS, WordCloud
import matplotlib.pyplot as plt
import datetime
import io
from PIL import Image 
import PIL 

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import ColorFormat, RGBColor

def app():

    file_path = "data-reportautomation/style.pptx"
    prs = Presentation(file_path)
    title_slide_layout = prs.slide_layouts[0]
    divider_slide_layout = prs.slide_layouts[3]
    content_slide_layout = prs.slide_layouts[2]
    text_slide_layout = prs.slide_layouts[5]
    end_slide_layout = prs.slide_layouts[19]

    st.header("Reporting Automation")
    st.write("Step 1: Enter your project name.")
    Project_name = st.text_input(label="Project Name:", value="Project Name")

    st.write("Step 2: choose the survey results file in below. ")
    data_file = st.file_uploader("Please select your .xls file:")
    #data_file = "data-reportautomation/data.xls"
    
    if data_file is not None:
        
        df_data = pd.read_excel(data_file, sheet_name="Data")
        labels = pd.read_excel(data_file, sheet_name="Labels")
        labels.columns = ["questions", "answers"]

        st.write("Project Data:")
        st.dataframe(df_data)

        #fragetexte = []

        df_data_encoded = pd.DataFrame()
        #### Labels
        df_label = pd.DataFrame()

        LABELS = {}
        questions = labels["questions"].str.split(" = ", 1, expand=True)
        answers = labels["answers"].str.split(" = ", 1, expand=True)

        df_label["q_key"] = questions[0]
        df_label["q_labels"] = questions[1]
        df_label["a_key"] = answers[0]
        df_label["a_labels"] = answers[1]

        for i, j in zip(questions[0], questions[1]):
            LABELS[i] = j
        #LABELS
        #df_label

        answer_keys = pd.DataFrame(columns=["Values"])

        st.write("Step 3: click 'calculate' to start processing the survey results.")
        if st.button("Calculate"): #to be replaced by file uploader - check is file is uploaded and start calc

            title_slide = prs.slides.add_slide(title_slide_layout)
            title_slide.shapes.title.text = Project_name  + "\n" + "Survey Results"
            title_slide.shapes[1].text = "Arne Lienemann \nConsumer Insights | Global Portfolio Management"

            divider_slide = prs.slides.add_slide(divider_slide_layout)
            divider_slide.shapes.title.text = "Results"

            for column in df_data.columns:
                if "f" in column:
                    
                    df_data = df_data[df_data != 999]
                                
                    absolute_frequencies = df_data[column].value_counts(normalize=False)
                    sum = absolute_frequencies.sum()

                    relative_frequencies = absolute_frequencies/sum
                    relative_frequencies.sort_index(inplace=True)
                    #relative_frequencies = pd.DataFrame(relative_frequencies)

                    if "u" in column:
                        if len(relative_frequencies) != 0:
                            text = df_data[column]
                            st.write(LABELS[column] + " (n = " + str(sum) + ")")
                            
                            x, y = np.ogrid[:400, :400]
                            mask = (x - 200) ** 2 + (y - 200) ** 2 > 175 ** 2
                            mask = 255 * mask.astype(int)
                            wc = WordCloud(font_path ="data-wordcloud/NotArial.ttf", stopwords=STOPWORDS, background_color="white", color_func=lambda *args, **kwargs: (0,144,214), width=400, height=400, mask=mask)
                            
                            text = ' '.join(df_data[column].fillna(''))
                            wordcloud = wc.generate(text) 

                            imageWC = wordcloud.to_image()
                            imageWC.save("data-reportautomation/saved_pic.png")
                            st.image(imageWC)

                            frequencies = WordCloud().process_text(text)
                            df_freq = pd.DataFrame(frequencies, index=["count"])
                            df_freq = df_freq.transpose()
                            df_freq_sorted = df_freq.sort_values("count", ascending=False) 
                            df_freq_sorted = df_freq_sorted.head(15)

                            content_slide = prs.slides.add_slide(text_slide_layout)
                            content_slide.shapes.title.text = "Open comments" 
                            
                            question_text = LABELS[column]
                            if "//" in question_text:
                                question_text = question_text.split("//")[1]

                            content_slide.shapes[0].text = question_text + " (n = " + str(sum) + ")." 
                            #content_slide.shapes[1].text.font.size = Pt(10.5) # TODO: has no attribute 'font'
                            #content_slide.shapes[2].text = "" #str(index + 3)
                            #content_slide.shapes[3].text = Project_name

                            chart_data = CategoryChartData()
                            chart_data.categories = df_freq_sorted.index.values
                            chart_data.add_series('Count', df_freq_sorted['count'])
                            graphic_frame = content_slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.72), Inches(2.09), Inches(5.94), Inches(4.65), chart_data)
                            graphic_frame.chart.series[0].format.fill.solid()
                            graphic_frame.chart.series[0].format.fill.fore_color.rgb = RGBColor(153, 174, 181)
                            graphic_frame.chart.value_axis.visible = False
                            graphic_frame.chart.category_axis.tick_labels.font.size = Pt(10.5)
                            graphic_frame.chart.category_axis.reverse_order = True
                            graphic_frame.chart.plots[0].has_data_labels = True
                            graphic_frame.chart.plots[0].data_labels.font.size = Pt(10.5)

                            pic = content_slide.shapes.add_picture("data-reportautomation/saved_pic.png", 0, 0)    
                            pic.left = int(prs.slide_width * 0.5)
                            pic.top = int((prs.slide_height - pic.height) * 0.5)   

                    elif "_" in column:
                        #st.write("Matrixfrage")
                        st.write(" ")
                        #st.write(column)

                    else:
                        if len(relative_frequencies) != 0:
                            
                            answer_keys = pd.DataFrame(columns=["Values"])
                            relative_frequencies.name = "Single choice"
                            index = df_label.index[df_label["q_key"] == column]
                            #df_label["q_labels"][max(index)] #max(index) weil teils "f3" zweimal vorkommt. Zuerst = [Single Punch], dann Fragetext
                            start_index = max(index) + 2
                            for answer_length in range(len(relative_frequencies)):
                                answer_keys.loc[answer_length] = df_label["a_labels"][start_index + answer_length]
                            relative_frequencies.index = answer_keys
                            
                            st.write( LABELS[column] + " (n = " + str(sum) + ")")
                            st.dataframe(relative_frequencies)

                            content_slide = prs.slides.add_slide(content_slide_layout)
                            
                            key_take_away = relative_frequencies.idxmax()[0]
                            content_slide.shapes.title.text = "Results" + "\n" + key_take_away

                            content_slide.shapes[1].text = LABELS[column] + " (n = " + str(sum) + ")."
                            #content_slide.shapes[2].text = "" #str(index + 3)
                            #content_slide.shapes[3].text = Project_name

                            chart_data = CategoryChartData()
                            chart_data.categories = answer_keys["Values"]
                            chart_data.add_series("", relative_frequencies)

                            if len(relative_frequencies) == 2:
                                graphic_frame = content_slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(0.72), Inches(1.35), Inches(9.89), Inches(5.39), chart_data) #TODO: Styling. Ringexplosion machbar?
                                graphic_frame.chart.plots[0].has_data_labels = True
                                graphic_frame.chart.plots[0].data_labels.font.size = Pt(10.5)
                                graphic_frame.chart.plots[0].data_labels.number_format = '0%'
                                graphic_frame.chart.has_legend = True
                                graphic_frame.chart.legend.font.size = Pt(10.5)
                            else:
                                graphic_frame = content_slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.73), Inches(1.35), Inches(9.89), Inches(5.39), chart_data)
                                graphic_frame.chart.series[0].format.fill.solid()
                                graphic_frame.chart.series[0].format.fill.fore_color.rgb = RGBColor(153, 174, 181)
                                graphic_frame.chart.value_axis.visible = False
                                graphic_frame.chart.category_axis.tick_labels.font.size = Pt(10.5)
                                graphic_frame.chart.plots[0].has_data_labels = True
                                graphic_frame.chart.plots[0].data_labels.font.size = Pt(10.5)
                                graphic_frame.chart.plots[0].data_labels.number_format = '0%'
            
            end_slide = prs.slides.add_slide(end_slide_layout)
            end_slide.shapes.title.text = "Thank you"
            end_slide.shapes[1].text ="Contact \nYour Name \nYour function \nyour@email.com"

            prs.save('data-reportautomation/newppt.pptx')
            st.write("Step 4: Download the final presentation as .pptx file.")
            #with open(file_path, 'rb') as my_file:
            with open('data-reportautomation/newppt.pptx', 'rb') as my_file:
                st.download_button(label = 'Download', data = my_file, file_name = Project_name + '.pptx') 


    #python -m streamlit run streamlit_reportautomation.py
