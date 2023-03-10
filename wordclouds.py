from configparser import Interpolation
from ctypes import util
from pickle import TRUE
#from turtle import color
from matplotlib import image
import streamlit as st
import numpy as np
import pandas as pd
from wordcloud import STOPWORDS, WordCloud
import matplotlib.pyplot as plt
import datetime
import io
from PIL import Image 
import PIL 

import streamlit as st
from langdetect import detect
from translate import Translator
from wordcloud import WordCloud

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import ColorFormat, RGBColor

def app():

    #uploaded_file = st.file_uploader("Choose a file")
    st.title("How to: Create a Wordcloud")
    question_text = "Word counts"
    #question_text = st.text_input(label="Enter your question here:", value="Question number 1:")
    text = st.text_area(label="Paste your text here:", value="abc")

    #### alternative: File-uploader
    #file = uploaded_file.read()
    #text = file.decode("utf-8")    

    # Counts the number of rows in text area. TODO: currently Buchstaben are counted
    line_count = 0
    for line in text:
        if line != "\n":
            line_count += 1
    #st.write(int(line_count))

    translate = ""
    #translate = st.checkbox("Automatically translate text into English:")

    if text is not "":

        x, y = np.ogrid[:400, :400]

        mask = (x - 200) ** 2 + (y - 200) ** 2 > 175 ** 2
        mask = 255 * mask.astype(int)

        #TODO: Add stopwords!
        #users_stopwords = st.text_input(label="Stopwords")
        #combined_stopwords = list(users_stopwords) + list(STOPWORDS)

        wc = WordCloud(font_path ="data-wordcloud/NotArial.ttf", stopwords=STOPWORDS, background_color="white", color_func=lambda *args, **kwargs: (0,144,214), width=400, height=400, mask=mask)

        # TODO: Add STOPWORDS!
        # TODO: Add translation via google translate!
        # TODO: Add Sennheiser Font       

        if translate is not "":

                translator = Translator(to_lang="en")
                translated_lines = []

                if translator:
                    for line in text.splitlines():
                        lang = detect(line)

                        if lang != 'en':
                            translated_lines.append(translator.translate(line))
                        else:
                            translated_lines.append(line)

                text = " ".join(translated_lines)
                wordcloud = wc.generate(text)

                imageWC = wordcloud.to_image()
                st.image(imageWC)

                #st.image(wordcloud, caption='Wordcloud', use_column_width=True)


        else:
            
            frequencies = WordCloud().process_text(text)
            df_freq = pd.DataFrame(frequencies, index=["count"])
            df_freq = df_freq.transpose()
            
            # Add: sort values by count
            df_freq_sorted = df_freq.sort_values("count", ascending=False) 

            df_freq_sorted = df_freq_sorted.head(15)
            #st.write(df_freq_sorted)

            #st.bar_chart(data=df_freq_sorted)

            wordcloud = wc.generate(text)

            imageWC = wordcloud.to_image()
            st.image(imageWC)

        imageWC.save("data-wordcloud/saved_pic.png")

    #Download-Section
#    if st.button('Create ppt'):
#
#        file_path = "data-wordcloud/template.pptx"
#        prs = Presentation(file_path)
#
#        slide = prs.slides[0]
#        shape = slide.shapes
#        shape[1].text = question_text + " (n = " + str(line_count) + ")"

#        chart_data = CategoryChartData()
#        chart_data.categories = df_freq_sorted.index.values
#        chart_data.add_series('Count', df_freq_sorted['count'])

#        chart = shape[2].chart
#        chart.replace_data(chart_data)

#        pic = slide.shapes.add_picture("data-wordcloud/saved_pic.png", 0, 0)    
#        pic.left = int(prs.slide_width * 0.5)
#        pic.top = int((prs.slide_height - pic.height) * 0.5)

#        prs.save('data-wordcloud/newppt.pptx')

        #with open(file_path, 'rb') as my_file:
#        with open('data-wordcloud/newppt.pptx', 'rb') as my_file:
#            st.download_button(label = 'Download', data = my_file, file_name = 'Wordloud.pptx') 

    #python -m streamlit run wordcloudingtostreamlit.py

    # mime = 
    #.xls      application/vnd.ms-excel
    #.xlsx     application/vnd.openxmlformats-officedocument.spreadsheetml.sheet
    #.ppt      application/vnd.ms-powerpoint
    #.pptx     application/vnd.openxmlformats-officedocument.presentationml.presentation
