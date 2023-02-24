from audioop import cross
from configparser import Interpolation
from ctypes import util
#from msilib.schema import ProgId
from pickle import TRUE
#from turtle import color
from matplotlib import image
import streamlit as st
import numpy as np
import pandas as pd
from wordcloud import WordCloud
import matplotlib.pyplot as plt
import datetime
import io
from PIL import Image 
import PIL 

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import ColorFormat, RGBColor

import operator
import csv

def app():

    data = None
    relatives = None

    st.header("Segmentation analysis")
    st.write("Step 1: Select the columns related to segmentation questions in your survey results excel sheet and paste them into a new excel file (.xls).")
    #st.write("Step 2: Click on the button below to start processing the segmentation analysis.")

    #data = pd.read_csv("data/example.csv", sep=";")

    uploaded_file = st.file_uploader("Choose a file")
    
    #if st.button(label="Read clipboard"):
     #   data = pd.read_clipboard(sep=";")

    if uploaded_file is not None:
        data = pd.read_excel(uploaded_file, names=["f_1","f_2","f_3", "f_4", "f_5","f_6","f_7", "f_8", "f_9","f_10"])
        key = pd.read_csv("segmentation-data/segmentation_key.csv", sep=";", index_col=0)
        #data = key[key.columns[1:6]]
        Segments = {}
        data["const"] = 1 #so that 1 is multiplied by the constant in key (11th row)
        segment = ""
        data_transposed = np.transpose(data)
        data_transposed = data_transposed.astype(float)
        #data_transposed

        st.write("Raw results: ")
        #for pid in list(data_transposed):

        for i in range(len(data)):
            pid = i
            #st.write(data_transposed[pid][int(q)])

            maxSum = -99
            segment = 1

            for columns in key:
                sum = 0
                for a, b in zip(data_transposed[pid], key[columns]):
                    sum += a*b
                    
                if (sum > maxSum):
                    #st.write("New MaxSum found!")
                    maxSum = sum
                    segment = columns
            Segments.update({pid: segment})

        #Session states are used so that streamlit doesn't forget values when re-running after second button click
        st.session_state.return_pd = pd.DataFrame(Segments.values(), index=Segments.keys(), columns=["Segment"])
        st.session_state.return_pd
        
        frequencies = pd.Series(st.session_state.return_pd["Segment"]).value_counts()
        st.session_state.relatives = frequencies/len(st.session_state.return_pd)
        st.session_state.relatives = st.session_state.relatives.round(2)
        st.session_state.relatives

        st.bar_chart(data = st.session_state.relatives)

    #TODO: Not working weith streamlit. ST has no access to the clipboard. Update!
    #if st.button("Copy raw data to clipboard"):
    #    st.session_state.return_pd["Segment"].to_clipboard()

    st.write("Step 2: Create PowerPoint slide based on calculated results:")
    if st.button('Create ppt'):
        file_path = "segmentation-data/template.pptx"
        prs = Presentation(file_path)

        slide = prs.slides[0]
        shapes = slide.shapes
        shapes[1].text = "How much do you agree or disagree to these statements describing your lifestyle? (n = " + str(len(st.session_state.return_pd)) + ")"
        
        chart_data = CategoryChartData()
        chart_data.categories = st.session_state.relatives.index.values
        chart_data.add_series('Segments', st.session_state.relatives)

        #Know your template! Using a template is a lot easier than creating a slide from scratch (even when working with placeholders)
        chart = shapes[4].chart
        chart.replace_data(chart_data)

        prs.save('segmentation-data/newppt.pptx')
        st.write("Step 3: PowerPoint slide created successfully! Click here to download: ")

        #with open(file_path, 'rb') as my_file:
        with open('segmentation-data/newppt.pptx', 'rb') as my_file:
            st.download_button(label = 'Download', data = my_file, file_name = 'Segmentation.pptx') 

#python -m streamlit run segmentation_analysis.py
